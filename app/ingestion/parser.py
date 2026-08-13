from io import BytesIO
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx"}


def parse_document(filename: str, data: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {suffix}")
    if suffix in {".txt", ".md"}:
        return data.decode("utf-8", errors="replace")
    if suffix == ".pdf":
        reader = PdfReader(BytesIO(data))
        return "\n\n".join(
            f"[page:{i}]\n{page.extract_text() or ''}"
            for i, page in enumerate(reader.pages, 1)
        )
    if suffix == ".docx":
        doc = Document(BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if suffix == ".pptx":
        prs = Presentation(BytesIO(data))
        blocks = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            blocks.append(f"[slide:{i}]\n" + "\n".join(texts))
        return "\n\n".join(blocks)
    workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
    blocks = []
    for ws in workbook.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                rows.append("\t".join(values))
        blocks.append(f"[sheet:{ws.title}]\n" + "\n".join(rows))
    return "\n\n".join(blocks)
